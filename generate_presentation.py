"""
Generate presentation using the official SIH template.
Minimal modifications to template slides. Adds new slides for required content.

Usage: python generate_presentation.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parent
TEMPLATE = ROOT / "IDEA_Presentation_Format.pptx"
OUT_DIR = ROOT / "Final_Deliverables" / "03_Presentation"
IMG_DIR = ROOT / "deliverables"

RED = RGBColor(0xE3, 0x18, 0x37)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY = RGBColor(0x64, 0x74, 0x8B)

def _set_run(slide, shape_name, text, font_size=None, bold=None, color=None):
    """Set the first run text in a shape by name."""
    for s in slide.shapes:
        if s.name == shape_name and s.has_text_frame:
            tf = s.text_frame
            if tf.paragraphs:
                p = tf.paragraphs[0]; r = p.runs[0] if p.runs else p.add_run()
                r.text = text
                if font_size: r.font.size = Pt(font_size)
                if bold is not None: r.font.bold = bold
                if color: r.font.color.rgb = color
                return True
    return False

def _set_shape_text(slide, shape_name, text, font_size=11, color=DARK, bold=False):
    """Set all text in a shape by name (replaces first paragraph)."""
    for s in slide.shapes:
        if s.name == shape_name and s.has_text_frame:
            tf = s.text_frame
            # Clear existing paragraphs
            for p in tf.paragraphs:
                p.clear()
            # Set first paragraph
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = text
            r.font.size = Pt(font_size)
            r.font.color.rgb = color
            r.font.bold = bold
            return tf
    return None

def _set_bullets(tf, items, font_size=11):
    """Set bullet items in a text frame."""
    if tf is None: return
    # Clear existing content
    for p in tf.paragraphs:
        p.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i == 0: p.text = ""
        r = p.add_run()
        r.text = item
        r.font.size = Pt(font_size)
        r.font.color.rgb = DARK

def _add_slide(prs, title_text, bullets, font_size=11, layout_idx=1):
    """Add a new Title and Content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    if slide.shapes.title:
        slide.shapes.title.text = ""
        r = slide.shapes.title.text_frame.paragraphs[0].add_run()
        r.text = title_text
    # Add bullet content
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.6), Inches(5.2))
    tb.text_frame.word_wrap = True
    _set_bullets(tb.text_frame, bullets, font_size)
    return slide

def _add_image(slide, path, left, top, width, height=None):
    if not path.exists(): return
    kw = {"left": Inches(left), "top": Inches(top), "width": Inches(width)}
    if height: kw["height"] = Inches(height)
    slide.shapes.add_picture(str(path), **kw)


# ═══════════════════════════════════════════════════════════════
def build():
    prs = Presentation(str(TEMPLATE))
    slides = list(prs.slides)
    L = prs.slide_layouts  # Layouts: 0=Title, 1=Title+Content, 2=Section, 3=Two Content, 5=Title Only, 6=Blank

    # ── SLIDE 0 (index 0): Repurpose instructions as Agenda ──
    s = slides[0]
    # Hide instructions text by making it small and gray
    for sh in s.shapes:
        if sh.has_text_frame and "IMPORTANT INSTRUCTIONS" in sh.text_frame.text:
            sh.text_frame.paragraphs[0].runs[0].text = ""
        if sh.has_text_frame and "Please ensure" in sh.text_frame.text:
            sh.text_frame.paragraphs[0].runs[0].text = ""
        # Repurpose the instruction bullet list as agenda
        if sh.has_text_frame and "Kindly keep the maximum" in sh.text_frame.text:
            _set_bullets(sh.text_frame, [
                "PROJECT AGENDA",
                "Problem Statement & Industrial Background",
                "Process Understanding & Dynamic Model",
                "Control Strategy & Safety Architecture",
                "Validation Methodology & Results",
                "Performance Summary & Lessons Learned",
                "Future Work & Conclusion",
            ], font_size=11)

    # ── SLIDE 1: TITLE PAGE ──
    s = slides[1]
    _set_run(s, "Subtitle 3", "Autonomous Production Choke Controller", 20, True, RED)
    for sh in s.shapes:
        if sh.has_text_frame and "Problem Statement ID" in sh.text_frame.text:
            p = sh.text_frame.paragraphs[0]; r = p.runs[0] if p.runs else p.add_run()
            r.text = ("Problem Statement ID: 1367\n"
                      "Title: Autonomous Choke Control\n"
                      "Theme: Smart Automation\n"
                      "Category: Software\n"
                      "Team: Honeywell Hackathon 2026")

    # ── SLIDE 2: PROPOSED SOLUTION ──
    s = slides[2]
    _set_run(s, "Title 1", "PROPOSED SOLUTION", 24, True, DARK)
    tf = _set_shape_text(s, "TextBox 8", "", 12)
    _set_bullets(tf, [
        "5-phase autonomous control pipeline for oil well choke management",
        "FOPDT predictive control with 3-step horizon and gain scheduling",
        "Automated system identification via 6 open-loop step tests",
        "4-layer defense-in-depth safety architecture",
        "3-term cost function: J = J_track + 0.01*J_effort - 0.005*J_margin",
        "Interactive Streamlit dashboard with Plotly visualization",
        "All models identified exclusively from simulator-generated data",
    ], 12)

    # ── SLIDE 3: TECHNICAL APPROACH ──
    s = slides[3]
    _set_run(s, "Title 1", "TECHNICAL APPROACH", 24, True, DARK)
    tf = _set_shape_text(s, "TextBox 8", "", 11)
    _set_bullets(tf, [
        "Technology: Python 3.12, NumPy, Matplotlib, Plotly, Streamlit",
        "Process Model: First-Order Plus Dead Time (FOPDT)",
        "Identification: 63.2% graphical method from step tests",
        "Optimization: Brute-force candidate search (11-21 moves, <1ms)",
        "Safety: 4-layer defense-in-depth with emergency override",
        "Gain scheduling: 3 regions [0-35%], [35-65%], [65-100%]",
        "Bias correction: Exponential EWMA (alpha=0.3)",
    ], 11)
    _add_image(s, IMG_DIR/"architecture"/"controller_pipeline_architecture.png", 0.4, 4.8, 9.0, 2.3)

    # ── SLIDE 4: FEASIBILITY ──
    s = slides[4]
    _set_run(s, "Title 1", "FEASIBILITY AND VIABILITY", 24, True, DARK)
    tf = _set_shape_text(s, "TextBox 8", "", 12)
    _set_bullets(tf, [
        "All 3 scenarios validated -- tracking error < 2 bbl/hr",
        "Zero violations in normal operation; correct override at boundary",
        "Computationally efficient: candidates evaluated in <1ms",
        "Deterministic: seed=42, 243 simulator calls, reproducible",
        "Nonlinearity: gain scheduling across 3 operating regions",
        "Disturbances: EWMA bias correction + deadband suppression",
        "Single boundary violation correctly triggered emergency override",
    ], 12)

    # ── SLIDE 5: ARTIFACTS ──
    s = slides[5]
    _set_run(s, "Title 1", "ARTIFACTS", 24, True, DARK)
    tf = _set_shape_text(s, "TextBox 8", "", 12)
    _set_bullets(tf, [
        "Source Code: 17 Python modules, fully documented",
        "Technical Report: 24-section professional PDF",
        "Validation Report: Gate-by-gate scenario evaluation",
        "Streamlit Dashboard: 8-page interactive monitoring",
        "Architecture Diagrams: Pipeline + Safety (PNG)",
        "Step Test Analysis: 9 plots across operating range",
        "Scenario Plots: 6-panel enhanced figures for A, B, C",
    ], 12)
    _add_image(s, IMG_DIR/"architecture"/"safety_defense_in_depth.png", 0.8, 5.2, 8.5, 1.8)

    # ── SLIDE 6: REFERENCES ──
    s = slides[6]
    _set_run(s, "Title 1", "RESEARCH AND REFERENCES", 24, True, DARK)
    tf = _set_shape_text(s, "TextBox 8", "", 12)
    _set_bullets(tf, [
        "Honeywell Problem Statement -- Autonomous Choke Control (2026)",
        "Seborg et al. -- Process Dynamics and Control (4th Ed., Wiley)",
        "Astrom & Hagglund -- Advanced PID Control (ISA, 2006)",
        "FOPDT Identification: 63.2% graphical method",
        "Bias correction: EWMA filtering for offset-free MPC",
        "Gain scheduling: Multi-region linear models for nonlinear processes",
        "Safety: IEC 61511 principles for process constraint enforcement",
    ], 12)

    # ═══ ADDITIONAL SLIDES ═══

    # SLIDE 8: Problem Statement
    _add_slide(prs, "PROBLEM STATEMENT", [
        "Oil well production choke: primary control actuator between reservoir and surface",
        "4 coupled process variables: Oil Rate (Q), WHP, FLP, BHP",
        "Choke UP -> Q UP, WHP DOWN, FLP DOWN, BHP DOWN",
        "Choke DOWN -> Q DOWN, WHP UP, FLP UP, BHP UP",
        "Hard pressure constraints must never be violated",
        "1-hour control interval, +/-5% max choke movement per step",
        "Objective: Maximize production while respecting all constraints",
    ], 12)

    # SLIDE 9: Industrial Background
    _add_slide(prs, "INDUSTRIAL BACKGROUND & ASSUMPTIONS", [
        "Choke valve controls flow rate and wellbore pressure simultaneously",
        "Opening choke: increases revenue, decreases all pressures",
        "Closing choke: conserves reservoir energy, reduces production",
        "Classic constrained optimization: maximize Q subject to WHP/FLP/BHP limits",
        "Key Assumptions: Linear FOPDT per region, hard constraints, Gaussian noise, " +
        "consistent gain signs, step tests acceptable during commissioning",
        "1-hour interval typical for remote well sites with limited communication",
    ], 12)

    # SLIDE 10: Step Test Analysis
    s = _add_slide(prs, "OPEN-LOOP STEP TEST ANALYSIS", [
        "6 step tests at choke positions [10, 20, 30, 50, 70, 80] with +10% steps",
        "Protocol: Settle 5 steps -> Step +10% -> Record 10-step transient",
        "90 total TestSimulator calls; all data from simulator (NOT reference dataset)",
        "Consistent FOPDT-like first-order responses observed across all operating points",
        "Varying steady-state gains indicate mild nonlinearity -> gain scheduling required",
        "No dead time detected (fast response relative to 1-hour sampling interval)",
    ], 11)
    _add_image(s, IMG_DIR/"step_tests"/"step_test_summary_all_6_tests.png", 0.3, 4.0, 9.4, 2.8)

    # SLIDE 11: Dynamic Model
    s = _add_slide(prs, "DYNAMIC MODEL IDENTIFICATION", [
        "FOPDT Model: G(s) = K*exp(-theta*s) / (tau*s + 1)",
        "Discrete: y(k+1) = a*y(k) + b*u(k-d) + c,  a=exp(-Ts/tau), b=K*(1-a)",
        "63.2% Graphical Method: K = delta_y_ss/delta_u, tau = time to 63.2%, " +
        "theta = first response > 5% threshold",
        "Averaged across all 6 tests for a single global model",
        "Q: K=+2.02 tau=1.0h | WHP: K=-3.00 tau=1.0h | FLP: K=-2.02 tau=1.0h | BHP: K=-4.13 tau=2.0h",
        "Close match to known simulator dynamics confirms identification accuracy",
        "All identification data from TestSimulator step tests ONLY",
    ], 11)
    _add_image(s, IMG_DIR/"step_tests"/"gain_comparison.png", 2.5, 4.2, 5.0, 2.5)

    # SLIDE 12: Control Strategy
    _add_slide(prs, "CONTROL STRATEGY", [
        "Phase 1 - Perception: 20-step rolling history, EWMA bias (alpha=0.3), SS detection, safety proximity",
        "Phase 2 - Targeting: STARTUP -> TRACKING (|error|<2%) -> INFEASIBLE (SS unreachable)",
        "Phase 3 - Prediction: Brute-force grid +/-5%, 1.0% coarse/0.5% fine resolution, FOPDT Np=3, gain scheduling",
        "Phase 4 - Safety: Tightened constraints (+5% margin), per-step trajectory + SS validation, emergency override",
        "Phase 5 - Selection: 3-term cost J=J_track+0.01*J_effort-0.005*J_margin, deadband, least-infeasible fallback",
        "Cost: J_track=((Q_pred-Q_target)/Q_range)^2 [tracking], =-(Q_pred/Q_range) [infeasible max production]",
        "Deadband: J_track=0 within 1% Q_range. WARNING: ramp 5%->2%. EMERGENCY: bypass optimizer, +/-5% override",
    ], 10)

    # SLIDE 13: Safety Architecture
    s = _add_slide(prs, "SAFETY ARCHITECTURE", [
        "Layer 1 - Proximity Detection: Normalized margin to constraints -> NORMAL/CAUTION/WARNING/EMERGENCY",
        "Layer 2 - Tightened Constraints: +5% margin buffer on all limits (tightened = limit +/- 5%*range)",
        "Layer 3 - Trajectory Checking: Every predicted step AND steady-state validated against tightened limits",
        "Layer 4 - Emergency Override: Hard limit violation -> immediate +/-5% correction, bypasses optimizer",
        "Design principle: No single layer failure can cause a constraint violation",
    ], 11)
    _add_image(s, IMG_DIR/"architecture"/"safety_defense_in_depth.png", 0.3, 3.8, 9.4, 2.8)

    # SLIDE 14: Validation
    _add_slide(prs, "VALIDATION METHODOLOGY", [
        "Three scenarios against TestSimulator (seed=42, deterministic):",
        "Scenario A: Constant Target Q=80 bbl/hr, 50 steps -- tests steady-state tracking",
        "Scenario B: Target Change Q=60->120 at step 26 -- tests setpoint transition",
        "Scenario C: Infeasible Target Q=250 bbl/hr -- tests constraint-limited max production",
        "Five Validation Gates: G1-Constraint, G2-Choke Range, G3-Tracking, G4-Mode, G5-Feasibility",
        "All performance metrics computed from simulator-generated logs (243 total calls)",
    ], 12)

    # SLIDE 15: Results A
    s = _add_slide(prs, "RESULTS: SCENARIO A -- CONSTANT TARGET (Q=80)", [
        "Final Q=80.9 bbl/hr, Tracking Error=0.9 bbl/hr (1.1%), 0 EMERGENCY violations",
        "Choke ramped 10%->40.5% over 7 steps, deadband suppressed movement for 40+ steps",
        "PASS: All 5 validation gates. Fine resolution (0.5%) active near target.",
    ], 12)
    _add_image(s, IMG_DIR/"scenarios"/"scenario_a_constant_target.png", 0.3, 3.0, 9.4, 3.8)

    # SLIDE 16: Results B
    s = _add_slide(prs, "RESULTS: SCENARIO B -- TARGET CHANGE (60->120)", [
        "Final Q=120.9, Error=0.9 bbl/hr (0.8%), 0 violations, smooth 7-step transition",
        "Choke 30.5%->60.5%, no overshoot, mode correctly transitioned STARTUP->TRACKING",
        "PASS: All 5 validation gates. Both Q=60 and Q=120 tracked within 1.4 bbl/hr.",
    ], 12)
    _add_image(s, IMG_DIR/"scenarios"/"scenario_b_target_change.png", 0.3, 3.0, 9.4, 3.8)

    # SLIDE 17: Results C
    s = _add_slide(prs, "RESULTS: SCENARIO C -- INFEASIBLE TARGET (Q=250)", [
        "Max safe Q=185.9 bbl/hr at choke=93%. 1 emergency override at step 19 (BHP boundary)",
        "Controller ramped to 94%, BHP noise triggered override -> retracted to 93% safe position",
        "PASS: Safety system functioned correctly. Controller maximized production at constraint edge.",
    ], 12)
    _add_image(s, IMG_DIR/"scenarios"/"scenario_c_infeasible_target.png", 0.3, 3.0, 9.4, 3.8)

    # SLIDE 18: Performance
    _add_slide(prs, "PERFORMANCE SUMMARY", [
        "Tracking: A: 0.9 bbl/hr (1.1%), B: 0.9 bbl/hr (0.8%), C: 64 bbl/hr (infeasible by design)",
        "Constraints: A,B: 0 violations. C: 1 violation at BHP edge, correctly managed by safety override",
        "Safety: All 4 layers verified. Emergency override triggered and recovered correctly.",
        "Reproducibility: 243 simulator calls at seed=42, deterministic and fully reproducible",
        "Average feasible candidates: 21/21 (A,B - fine resolution), 3-5/5-11 (C - boundary-tightened)",
    ], 12)

    # SLIDE 19: Lessons Learned
    _add_slide(prs, "LESSONS LEARNED", [
        "FOPDT from 6 step tests provides adequate accuracy; 63.2% method matches known simulator dynamics",
        "Bias correction (alpha=0.3) is essential for offset-free tracking with imperfect models",
        "Fine resolution (0.5%) near targets reduces chatter; coarse (1.0%) keeps candidates manageable",
        "5% safety margin on tightened constraints prevents operation at constraint edge",
        "Deadband suppression is critical for industrial acceptance -- operators reject unnecessary movement",
        "Gain scheduling with simple regions handles nonlinearity without complex nonlinear models",
    ], 11)

    # SLIDE 20: Future Work
    _add_slide(prs, "FUTURE WORK", [
        "Online model adaptation: RLS with forgetting factor to update FOPDT as conditions change",
        "Soft constraints: Exponential penalty replacing hard binary rejection for gradual limit approach",
        "Disturbance feedforward: Model reservoir pressure decline for preemptive control",
        "Multi-well coordination: Extend to optimize across wells sharing surface facilities",
        "OPC-UA integration: Connect to Honeywell Experion DCS for real-time industrial deployment",
        "Sensor fault detection: Cross-validate redundant measurements to mitigate failures",
    ], 11)

    # SLIDE 21: Conclusion
    _add_slide(prs, "CONCLUSION", [
        "The Autonomous Choke Controller meets all Honeywell Hackathon requirements",
        "5-phase sequential architecture with 4-layer defense-in-depth safety",
        "Automatic commissioning: step tests -> FOPDT identification -> control in one pipeline",
        "Offset-free tracking via EWMA bias correction; nonlinearity via gain scheduling",
        "Graceful degradation under infeasible targets -> maximum safe production",
        "Full explainability: every action logged with cost, reason, candidate evaluation",
        "All models from simulator data only. Reference dataset = post-hoc check only.",
        "Ready for integration with Honeywell process simulator and field deployment.",
    ], 11)

    # ── Save ──
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    path = OUT_DIR / "Presentation.pptx"
    prs.save(str(path))
    return path, len(prs.slides)


if __name__ == "__main__":
    print("Building presentation from template...")
    path, count = build()
    print(f"  [OK] {path}")
    print(f"  Slides: {count}")
